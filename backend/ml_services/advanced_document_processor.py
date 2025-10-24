import os
import json
import re
import PyPDF2
import io

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

class AdvancedDocumentProcessor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self.process_pdf,
            '.docx': self.process_docx,
            '.pptx': self.process_pptx,
            '.xlsx': self.process_excel,
            '.csv': self.process_csv
        }
    
    def process_document(self, file_path):
        """Process document and extract structured data"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported format: {ext}")
        
        processor = self.supported_formats[ext]
        raw_data = processor(file_path)
        
        return self.extract_structured_data(raw_data)
    
    def process_pdf(self, file_path):
        """Enhanced PDF processing with OCR"""
        text_content = []
        images = []
        
        if fitz:  # Use PyMuPDF if available
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                text_content.append(text)
                
                # OCR if available
                if pytesseract and Image:
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n < 5:
                                img_data = pix.tobytes("png")
                                img_pil = Image.open(io.BytesIO(img_data))
                                ocr_text = pytesseract.image_to_string(img_pil)
                                if ocr_text.strip():
                                    images.append({
                                        'page': page_num + 1,
                                        'text': ocr_text,
                                        'type': 'chart_or_image'
                                    })
                            pix = None
                        except:
                            continue
            doc.close()
        else:  # Fallback to PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
        
        return {
            'text': '\n'.join(text_content),
            'images': images,
            'page_count': len(text_content)
        }
    
    def process_docx(self, file_path):
        """Process Word documents"""
        if not Document:
            return {'error': 'python-docx not installed'}
            
        try:
            doc = Document(file_path)
            text_content = []
            tables = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                'text': '\n'.join(text_content),
                'tables': tables
            }
        except Exception as e:
            return {'error': str(e)}
    
    def process_pptx(self, file_path):
        """Process PowerPoint presentations"""
        if not Presentation:
            return {'error': 'python-pptx not installed'}
            
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                slides_content.append({
                    'slide_number': slide_num + 1,
                    'content': '\n'.join(slide_text)
                })
            
            return {
                'slides': slides_content,
                'total_slides': len(slides_content)
            }
        except Exception as e:
            return {'error': str(e)}
    
    def process_excel(self, file_path):
        """Process Excel files"""
        if not pd:
            return {'error': 'pandas not installed'}
            
        try:
            excel_data = pd.read_excel(file_path, sheet_name=None)
            processed_sheets = {}
            for sheet_name, df in excel_data.items():
                processed_sheets[sheet_name] = {
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'shape': df.shape
                }
            return processed_sheets
        except Exception as e:
            return {'error': str(e)}
    
    def process_csv(self, file_path):
        """Process CSV files"""
        if not pd:
            return {'error': 'pandas not installed'}
            
        try:
            df = pd.read_csv(file_path)
            return {
                'data': df.to_dict('records'),
                'columns': df.columns.tolist(),
                'shape': df.shape
            }
        except Exception as e:
            return {'error': str(e)}
    
    def extract_structured_data(self, raw_data):
        """Extract structured business data from raw content"""
        if isinstance(raw_data, dict) and 'text' in raw_data:
            text = raw_data['text']
        else:
            text = str(raw_data)
        
        return {
            'company_info': self.extract_company_info(text),
            'financial_data': self.extract_financial_data(text),
            'team_info': self.extract_team_info(text),
            'market_data': self.extract_market_data(text),
            'product_info': self.extract_product_info(text),
            'raw_data': raw_data
        }
    
    def extract_company_info(self, text):
        """Extract company information"""
        company_info = {}
        
        # Company name patterns
        name_patterns = [
            r'Company:\s*([^\n]+)',
            r'Startup:\s*([^\n]+)',
            r'^([A-Z][a-zA-Z\s]+)(?:\s+Inc\.|\s+LLC|\s+Corp\.)?',
        ]
        
        for pattern in name_patterns:
            match = re.search(pattern, text, re.MULTILINE | re.IGNORECASE)
            if match:
                company_info['name'] = match.group(1).strip()
                break
        
        # Founded year
        year_match = re.search(r'(?:founded|established|started).*?(\d{4})', text, re.IGNORECASE)
        if year_match:
            company_info['founded_year'] = year_match.group(1)
        
        # Location
        location_patterns = [
            r'(?:based|located|headquarters?).*?in\s+([^,\n]+(?:,\s*[A-Z]{2})?)',
            r'([A-Z][a-z]+,\s*[A-Z]{2})',
        ]
        
        for pattern in location_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                company_info['location'] = match.group(1).strip()
                break
        
        return company_info
    
    def extract_financial_data(self, text):
        """Extract financial information"""
        financial_data = {}
        
        # Revenue patterns
        revenue_patterns = [
            r'revenue.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
            r'\$([0-9,]+(?:\.[0-9]+)?[KMB]?).*?revenue',
            r'ARR.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
        ]
        
        for pattern in revenue_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                financial_data['revenue'] = match.group(1)
                break
        
        # Funding patterns
        funding_patterns = [
            r'raised.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
            r'funding.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
            r'seeking.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
        ]
        
        for pattern in funding_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                financial_data['funding'] = match.group(1)
                break
        
        # Growth rate
        growth_match = re.search(r'(\d+)%.*?growth', text, re.IGNORECASE)
        if growth_match:
            financial_data['growth_rate'] = growth_match.group(1) + '%'
        
        return financial_data
    
    def extract_team_info(self, text):
        """Extract team information"""
        team_info = {}
        
        # Founder patterns
        founder_patterns = [
            r'(?:founder|CEO|CTO|co-founder).*?([A-Z][a-z]+\s+[A-Z][a-z]+)',
            r'([A-Z][a-z]+\s+[A-Z][a-z]+).*?(?:founder|CEO|CTO)',
        ]
        
        founders = []
        for pattern in founder_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            founders.extend(matches)
        
        team_info['founders'] = list(set(founders))  # Remove duplicates
        
        # Team size
        team_size_match = re.search(r'team.*?(\d+).*?(?:people|employees|members)', text, re.IGNORECASE)
        if team_size_match:
            team_info['team_size'] = team_size_match.group(1)
        
        return team_info
    
    def extract_market_data(self, text):
        """Extract market information"""
        market_data = {}
        
        # Market size patterns
        market_patterns = [
            r'market.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
            r'TAM.*?\$([0-9,]+(?:\.[0-9]+)?[KMB]?)',
            r'\$([0-9,]+(?:\.[0-9]+)?[KMB]?).*?market',
        ]
        
        for pattern in market_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                market_data['market_size'] = match.group(1)
                break
        
        # Industry/sector
        sector_keywords = ['fintech', 'healthtech', 'edtech', 'saas', 'e-commerce', 'ai', 'blockchain']
        for keyword in sector_keywords:
            if keyword.lower() in text.lower():
                market_data['sector'] = keyword.title()
                break
        
        return market_data
    
    def extract_product_info(self, text):
        """Extract product information"""
        product_info = {}
        
        # Product type
        product_keywords = ['platform', 'app', 'software', 'service', 'solution', 'tool']
        for keyword in product_keywords:
            if keyword.lower() in text.lower():
                product_info['type'] = keyword.title()
                break
        
        # Key features (simple extraction)
        feature_patterns = [
            r'features?.*?:([^.]+)',
            r'capabilities?.*?:([^.]+)',
        ]
        
        for pattern in feature_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                product_info['features'] = match.group(1).strip()
                break
        
        return product_info